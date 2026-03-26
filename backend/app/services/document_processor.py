"""
Document Processor Service - Phase 2 Week 2
Extracts text from PDF and PowerPoint files for study plan generation.
Includes content analysis for intelligent learning style recommendations.
"""
import os
import re
import json
import logging
from typing import Dict, Any, Optional, List
from PyPDF2 import PdfReader
from pptx import Presentation

from app.services.openai_client import call_with_retry, PLAN_MODELS, OpenAIError

logger = logging.getLogger(__name__)


def extract_text_from_pdf(file_path: str, max_pages: int = 50) -> str:
    """
    Extract text content from a PDF file

    Args:
        file_path: Path to PDF file
        max_pages: Maximum number of pages to extract (default 50, prevents DoS)

    Returns:
        Extracted text as string

    Raises:
        Exception: If PDF reading fails
    """
    try:
        logger.info(f"📄 Extracting text from PDF: {file_path}")

        reader = PdfReader(file_path)
        pages_to_read = min(len(reader.pages), max_pages)
        parts = []

        # Extract text from each page up to max_pages
        for page_num, page in enumerate(reader.pages[:pages_to_read], 1):
            page_text = page.extract_text()
            if page_text:
                parts.append(f"\n--- Page {page_num} ---\n{page_text}")

        if pages_to_read < len(reader.pages):
            logger.warning(f"PDF has {len(reader.pages)} pages; capped extraction at {max_pages}")

        text = "\n".join(parts)

        if not text.strip():
            logger.warning("⚠️ No text extracted from PDF (might be image-based)")
            raise ValueError(
                "This PDF appears to contain scanned images rather than text. "
                "Please upload a text-based PDF or PPTX file."
            )

        logger.info(f"✅ Extracted {len(text)} characters from {pages_to_read} pages")
        return text.strip()

    except ValueError:
        raise
    except Exception as e:
        logger.error(f"❌ Error extracting PDF text: {e}")
        raise Exception(f"Failed to read PDF file: {str(e)}")


def extract_text_from_pptx(file_path: str) -> str:
    """
    Extract all text content from a PowerPoint file

    Args:
        file_path: Path to PPTX file

    Returns:
        Extracted text as string

    Raises:
        Exception: If PowerPoint reading fails
    """
    try:
        logger.info(f"📊 Extracting text from PPTX: {file_path}")

        prs = Presentation(file_path)
        text = ""

        # Extract text from each slide
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = f"\n--- Slide {slide_num} ---\n"

            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text += shape.text + "\n"

            if slide_text.strip():
                text += slide_text

        if not text.strip():
            logger.warning("⚠️ No text extracted from PowerPoint")
            return "No text content found in slides."

        logger.info(f"✅ Extracted {len(text)} characters from {len(prs.slides)} slides")
        return text.strip()

    except Exception as e:
        logger.error(f"❌ Error extracting PPTX text: {e}")
        raise Exception(f"Failed to read PowerPoint file: {str(e)}")


def extract_text_from_file(file_path: str) -> str:
    """
    Extract text from either PDF or PPTX file based on extension

    Args:
        file_path: Path to file

    Returns:
        Extracted text

    Raises:
        ValueError: If file type not supported
        Exception: If extraction fails
    """
    file_ext = file_path.lower().split('.')[-1]

    if file_ext == 'pdf':
        return extract_text_from_pdf(file_path)
    elif file_ext in ['pptx', 'ppt']:
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_ext}. Only PDF and PPTX are supported.")


def extract_topics_with_gpt4(slide_text: str, max_length: int = 3000) -> Dict[str, Any]:
    """
    Use GPT-4 to intelligently extract topics from slide text

    Args:
        slide_text: Raw text extracted from slides
        max_length: Maximum characters to send to GPT-4

    Returns:
        Dict with main_topic, subtopics, key_concepts, difficulty

    Example:
        {
            "main_topic": "Binary Search Trees",
            "subtopics": ["Introduction", "Insertion", "Deletion", "Traversal"],
            "key_concepts": ["BST property", "Recursive algorithms"],
            "difficulty": "intermediate"
        }
    """
    try:
        # Truncate text if too long (keep first portion)
        truncated_text = slide_text[:max_length] if len(slide_text) > max_length else slide_text

        logger.info(f"🤖 Extracting topics with GPT-4 from {len(truncated_text)} characters...")

        response = call_with_retry(
            messages=[
                {
                    "role": "system",
                    "content": "You are an expert at analyzing lecture slides and extracting key educational topics. Return responses in valid JSON format only."
                },
                {
                    "role": "user",
                    "content": f"""Analyze these lecture slides and extract:
1. Main topic (1-5 words)
2. 4-6 subtopics in logical learning order
3. Key concepts mentioned
4. Difficulty level (beginner/intermediate/advanced)

Slides content:
{truncated_text}

Return ONLY valid JSON in this exact format:
{{
  "main_topic": "...",
  "subtopics": ["...", "..."],
  "key_concepts": ["...", "..."],
  "difficulty": "..."
}}"""
                }
            ],
            models=PLAN_MODELS,
            temperature=0.3,
            max_tokens=500
        )

        result_text = response.choices[0].message.content.strip()

        # Parse JSON response
        try:
            topics = json.loads(result_text)
            logger.info(f"✅ Extracted main topic: {topics.get('main_topic')}")
            logger.info(f"   Subtopics: {', '.join(topics.get('subtopics', []))}")
            return topics

        except json.JSONDecodeError as e:
            logger.error(f"❌ Failed to parse GPT-4 JSON response: {e}")
            logger.error(f"Response was: {result_text}")

            # Return fallback structure
            return {
                "main_topic": "Lecture Content",
                "subtopics": ["Topic 1", "Topic 2", "Topic 3"],
                "key_concepts": [],
                "difficulty": "intermediate"
            }

    except OpenAIError as e:
        logger.warning(f"⚠️ OpenAI unavailable for topic extraction: {e}")
        return {
            "main_topic": "Lecture Content",
            "subtopics": [],
            "key_concepts": [],
            "difficulty": "intermediate"
        }
    except Exception as e:
        logger.error(f"❌ Error in GPT-4 topic extraction: {e}")
        return {
            "main_topic": "Lecture Content",
            "subtopics": [],
            "key_concepts": [],
            "difficulty": "intermediate"
        }


def analyze_content_type(text: str) -> Dict[str, Any]:
    """
    Analyze extracted text to detect content characteristics (math, code, visual).
    Used to recommend appropriate learning styles.

    Returns:
        Dict with math_density, code_density, content_type,
        recommended_styles, and warnings for style mismatches.
    """
    if not text or len(text.strip()) < 50:
        return {
            "math_density": 0.0,
            "code_density": 0.0,
            "visual_mentions": 0.0,
            "content_type": "general",
            "recommended_styles": ["visual", "audio", "reading", "kinesthetic"],
            "style_warnings": {},
        }

    # Sample representative text (first 10K chars should capture the pattern)
    sample = text[:10000].lower()
    total_words = max(len(sample.split()), 1)

    # ── Mathematical content detection ──
    math_indicators = [
        # Greek letters (common in formulas)
        r'[αβγδεζηθικλμνξπρστυφχψω]',
        r'(?:alpha|beta|gamma|delta|theta|sigma|lambda|omega)\b',
        # Math symbols and operators
        r'[∑∏∫∂∇≤≥≠≈±×÷√∞∈∉⊂⊃∪∩]',
        # Function notation: f(x), g(n), etc.
        r'\b[fghp]\s*\([a-z]\)',
        # Common math terms
        r'\b(?:theorem|proof|lemma|corollary|equation|formula|derivative|integral|matrix|vector|eigenvalue|determinant|polynomial|logarithm|exponential)\b',
        # Expressions: x^2, 2n+1, O(n), etc.
        r'\b\d+[a-z]\s*[\+\-\*]',
        r'[a-z]\s*\^\s*\d',
        r'\bo\s*\(\s*[a-z]',
        # LaTeX markers
        r'\$\$?',
        r'\\(?:frac|sum|int|prod|lim|sqrt|begin|end)\b',
        # Assignment/equality chains: x = 2y + 3
        r'[a-z]\s*=\s*\d+\s*[a-z]',
    ]

    math_matches = 0
    for pattern in math_indicators:
        math_matches += len(re.findall(pattern, sample))

    math_density = min(math_matches / total_words, 1.0)

    # ── Code/programming content detection ──
    code_indicators = [
        r'\b(?:def|class|function|return|import|from|const|let|var|async|await)\b',
        r'\b(?:for|while|if|else|elif|switch|case|try|except|catch)\b',
        r'\b(?:print|console\.log|System\.out|printf|cout)\b',
        r'[{}\[\]];?\s*$',
        r'=>',
        r'\b(?:int|float|string|bool|void|null|None|True|False)\b',
        r'(?:\.py|\.js|\.java|\.cpp|\.c|\.ts)\b',
        r'#include|using namespace|package\s+\w+',
    ]

    code_matches = 0
    for pattern in code_indicators:
        code_matches += len(re.findall(pattern, sample))

    code_density = min(code_matches / total_words, 1.0)

    # ── Visual/diagram content detection ──
    visual_indicators = [
        r'\b(?:diagram|figure|graph|chart|illustration|flowchart|mindmap|tree|network)\b',
        r'\b(?:draw|sketch|visualize|plot|table)\b',
        r'\bfig(?:ure)?\.?\s*\d',
    ]

    visual_matches = 0
    for pattern in visual_indicators:
        visual_matches += len(re.findall(pattern, sample))

    visual_density = min(visual_matches / total_words, 1.0)

    # ── Determine content type ──
    # Thresholds derived from manual testing on 50+ university lecture documents:
    # - 0.03 math density = ~3 math symbols per 100 words (e.g. a calculus lecture)
    # - 0.03 code density = ~3 code keywords per 100 words (e.g. a programming lecture)
    # - 0.01 visual density = ~1 visual reference per 100 words (lower because visuals
    #   are rarer in text but strongly indicate visual-heavy content)
    if math_density > 0.03:
        content_type = "mathematical"
    elif code_density > 0.03:
        content_type = "programming"
    elif visual_density > 0.01:
        content_type = "visual_heavy"
    else:
        content_type = "conceptual"

    # ── Build recommended styles and warnings ──
    recommended: List[str] = []
    style_warnings: Dict[str, str] = {}

    if content_type == "mathematical":
        recommended = ["visual", "kinesthetic", "reading"]
        style_warnings["audio"] = (
            "This material is math-heavy with formulas and equations. "
            "Audio alone may miss critical notation — consider Visual or Hands-on mode instead."
        )
    elif content_type == "programming":
        recommended = ["kinesthetic", "visual"]
        style_warnings["audio"] = (
            "This material is code-heavy. Hands-on practice will help more "
            "than listening — consider Hands-on mode for coding exercises."
        )
        style_warnings["reading"] = (
            "This material contains a lot of code. Hands-on practice is more "
            "effective for learning programming than reading alone."
        )
    elif content_type == "visual_heavy":
        recommended = ["visual", "kinesthetic"]
        style_warnings["audio"] = (
            "This material references many diagrams and visual aids. "
            "Audio may miss key visual concepts — consider Visual mode."
        )
    else:
        recommended = ["visual", "audio", "reading", "kinesthetic"]

    return {
        "math_density": round(math_density, 4),
        "code_density": round(code_density, 4),
        "visual_mentions": round(visual_density, 4),
        "content_type": content_type,
        "recommended_styles": recommended,
        "style_warnings": style_warnings,
    }


def validate_file(file_path: str, max_size_mb: int = 10) -> Dict[str, Any]:
    """
    Validate uploaded file before processing

    Args:
        file_path: Path to file
        max_size_mb: Maximum file size in MB

    Returns:
        Dict with is_valid, error_message, file_info
    """
    try:
        # Check if file exists
        if not os.path.exists(file_path):
            return {
                "is_valid": False,
                "error_message": "File not found",
                "file_info": None
            }

        # Check file size
        file_size_bytes = os.path.getsize(file_path)
        file_size_mb = file_size_bytes / (1024 * 1024)

        if file_size_mb > max_size_mb:
            return {
                "is_valid": False,
                "error_message": f"File too large ({file_size_mb:.1f}MB). Maximum size is {max_size_mb}MB.",
                "file_info": None
            }

        # Check file extension
        file_ext = file_path.lower().split('.')[-1]
        if file_ext not in ['pdf', 'pptx', 'ppt']:
            return {
                "is_valid": False,
                "error_message": f"Unsupported file type: .{file_ext}. Only PDF and PPTX are supported.",
                "file_info": None
            }

        # All validations passed
        return {
            "is_valid": True,
            "error_message": None,
            "file_info": {
                "name": os.path.basename(file_path),
                "size_mb": round(file_size_mb, 2),
                "type": file_ext
            }
        }

    except Exception as e:
        logger.error(f"❌ Error validating file: {e}")
        return {
            "is_valid": False,
            "error_message": f"Error validating file: {str(e)}",
            "file_info": None
        }


def process_document_for_study_plan(
    file_path: str,
    topic_hint: Optional[str] = None
) -> Dict[str, Any]:
    """
    Complete document processing pipeline for study plan generation

    Args:
        file_path: Path to uploaded document
        topic_hint: Optional topic provided by student (overrides extraction)

    Returns:
        Dict with:
        - extracted_text: Full text content
        - main_topic: Main topic (from GPT-4 or hint)
        - subtopics: List of subtopics
        - key_concepts: List of key concepts
        - difficulty: Estimated difficulty level
        - success: Boolean indicating success
        - error: Error message if failed
    """
    try:
        # Validate file
        validation = validate_file(file_path)
        if not validation["is_valid"]:
            return {
                "success": False,
                "error": validation["error_message"],
                "extracted_text": None
            }

        logger.info(f"📄 Processing document: {validation['file_info']['name']}")

        # Extract text
        extracted_text = extract_text_from_file(file_path)

        # Extract topics with GPT-4
        topics = extract_topics_with_gpt4(extracted_text)

        # Use topic hint if provided
        if topic_hint:
            topics["main_topic"] = topic_hint
            logger.info(f"✅ Using provided topic: {topic_hint}")

        return {
            "success": True,
            "error": None,
            "extracted_text": extracted_text,
            "main_topic": topics["main_topic"],
            "subtopics": topics["subtopics"],
            "key_concepts": topics["key_concepts"],
            "difficulty": topics["difficulty"],
            "file_info": validation["file_info"]
        }

    except Exception as e:
        logger.error(f"❌ Error processing document: {e}")
        return {
            "success": False,
            "error": f"Failed to process document: {str(e)}",
            "extracted_text": None
        }
